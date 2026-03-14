#!/usr/bin/env python3
"""
Y2K Global AI Portfolio Generator
Creates a professional PPTX presentation showcasing AI projects
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

def hex_to_rgb(hex_color):
    """Convert hex color to RGB tuple"""
    hex_color = hex_color.lstrip('#')
    return tuple(int(hex_color[i:i+2], 16) for i in (0, 2, 4))

def create_ai_portfolio():
    # Create presentation with 16:9 aspect ratio
    prs = Presentation()
    prs.slide_width = Inches(13.33)  # 16:9 widescreen
    prs.slide_height = Inches(7.5)
    
    # Define colors
    bg_color = RGBColor(*hex_to_rgb('#0a0f1c'))  # Dark navy/charcoal
    text_color = RGBColor(255, 255, 255)  # White
    accent_color = RGBColor(*hex_to_rgb('#00d4ff'))  # Electric blue
    
    # Slide 1: Cover
    slide_layout = prs.slide_layouts[6]  # Blank layout
    slide = prs.slides.add_slide(slide_layout)
    
    # Set background
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = bg_color
    
    # Main title
    title_box = slide.shapes.add_textbox(Inches(1), Inches(1.5), Inches(11.33), Inches(1.5))
    title_frame = title_box.text_frame
    title_frame.text = "AI PORTFOLIO"
    title_para = title_frame.paragraphs[0]
    title_para.alignment = PP_ALIGN.CENTER
    title_para.font.size = Pt(48)
    title_para.font.bold = True
    title_para.font.color.rgb = accent_color
    
    # Subtitle
    subtitle_box = slide.shapes.add_textbox(Inches(1), Inches(3), Inches(11.33), Inches(0.8))
    subtitle_frame = subtitle_box.text_frame
    subtitle_frame.text = "AI Systems Architecture & Consulting"
    subtitle_para = subtitle_frame.paragraphs[0]
    subtitle_para.alignment = PP_ALIGN.CENTER
    subtitle_para.font.size = Pt(24)
    subtitle_para.font.color.rgb = text_color
    
    # Company info
    company_box = slide.shapes.add_textbox(Inches(1), Inches(4.2), Inches(11.33), Inches(2))
    company_frame = company_box.text_frame
    company_text = """Y2K Global OÜ + KISS IT Solutions e.U.
Gergo Kiss, MSc — Managing Director & AI Systems Architect

+43 1 442 20 143 | office@y2k.global | www.y2k.global"""
    company_frame.text = company_text
    company_para = company_frame.paragraphs[0]
    company_para.alignment = PP_ALIGN.CENTER
    company_para.font.size = Pt(16)
    company_para.font.color.rgb = text_color
    
    for para in company_frame.paragraphs[1:]:
        para.alignment = PP_ALIGN.CENTER
        para.font.size = Pt(16)
        para.font.color.rgb = text_color

    def create_project_slide(title, description, tech_stack, category):
        slide = prs.slides.add_slide(slide_layout)
        
        # Set background
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = bg_color
        
        # Category label
        cat_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(3), Inches(0.5))
        cat_frame = cat_box.text_frame
        cat_frame.text = category
        cat_para = cat_frame.paragraphs[0]
        cat_para.font.size = Pt(12)
        cat_para.font.color.rgb = accent_color
        cat_para.font.bold = True
        
        # Project title
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(1), Inches(12.33), Inches(1))
        title_frame = title_box.text_frame
        title_frame.text = title
        title_para = title_frame.paragraphs[0]
        title_para.font.size = Pt(32)
        title_para.font.bold = True
        title_para.font.color.rgb = text_color
        
        # Description
        desc_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.2), Inches(12.33), Inches(3))
        desc_frame = desc_box.text_frame
        desc_frame.text = description
        desc_para = desc_frame.paragraphs[0]
        desc_para.font.size = Pt(14)
        desc_para.font.color.rgb = text_color
        desc_frame.word_wrap = True
        
        # Tech stack
        tech_box = slide.shapes.add_textbox(Inches(0.5), Inches(5.8), Inches(12.33), Inches(1.2))
        tech_frame = tech_box.text_frame
        tech_frame.text = f"Technologies: {tech_stack}"
        tech_para = tech_frame.paragraphs[0]
        tech_para.font.size = Pt(12)
        tech_para.font.color.rgb = accent_color
        tech_para.font.italic = True

    # Slide 2: Funding Cockpit
    create_project_slide(
        "Funding Cockpit",
        """AI-powered funding opportunity platform that revolutionizes grant application management. 
        
The platform automatically discovers relevant funding opportunities, evaluates their fit with your projects, generates high-quality application drafts, and tracks the entire application lifecycle from submission to outcome.

Currently live and serving clients at funding.kiss-it.io, the platform significantly reduces the time and effort required for successful grant applications while improving success rates through AI-driven optimization.""",
        "Python, FastAPI, PostgreSQL, AI/ML, Email Integration, Web Scraping",
        "FUNDING & GRANTS"
    )
    
    # Slide 3: Cortex
    create_project_slide(
        "Cortex",
        """Production-grade CPU-based Large Language Model inference platform designed for enterprise deployment.

Features an OpenAI-compatible API for seamless integration, intelligent multi-model routing for optimal performance, and a comprehensive admin dashboard for monitoring and management.

Built to scale horizontally while maintaining consistent performance, Cortex enables organizations to deploy AI capabilities without the complexity of GPU infrastructure management.""",
        "Python, LLM, API Gateway, Docker",
        "AI INFRASTRUCTURE"
    )
    
    # Slide 4: DAIMOND Platform
    create_project_slide(
        "DAIMOND Platform",
        """Private AI platform for document intelligence that transforms how organizations interact with their knowledge base.

The platform automatically processes uploaded documents through intelligent chunking and embedding generation, enabling powerful vector-based search through ChromaDB. Users can engage in natural language conversations with their documents using RAG (Retrieval-Augmented Generation) with full citation tracking.

Designed for enterprise privacy and security, ensuring sensitive documents never leave your infrastructure.""",
        "Python, ChromaDB, RAG, Embeddings, NLP",
        "DOCUMENT AI"
    )
    
    # Slide 5: FIS AI
    create_project_slide(
        "FIS AI",
        """AI companion for Fleet Information System that modernizes fleet management through intelligent automation.

Enables natural language querying of fleet data, advanced vision-based document processing for automated data extraction, and implements a ReAct-style tool-calling agent architecture with comprehensive safety guardrails.

Seamlessly integrates with existing enterprise systems while providing an intuitive interface for fleet operators and managers.""",
        "Python, Vision AI, NLP, Tool-Calling, Enterprise Integration",
        "ENTERPRISE AI"
    )
    
    # Slide 6: Master Project ML
    create_project_slide(
        "Master Project ML",
        """Machine learning powered bid creation system for logistics optimization that transforms transportation planning.

Leverages advanced ML algorithms to analyze historical data, route optimization, cost factors, and market conditions to generate competitive and profitable bids automatically.

Significantly reduces bid preparation time while improving win rates through data-driven decision making and optimization algorithms.""",
        "Python, ML, Scikit-learn, Optimization",
        "LOGISTICS AI"
    )
    
    # Slide 7: Supporting Projects
    create_project_slide(
        "Supporting Projects",
        """A diverse portfolio of AI and automation solutions spanning multiple domains:

• Master Thesis (2020): Machine Learning applications in Accounting automation
• Bookkeeping Automation: SaaS platform with AI-powered invoice extraction and processing
• OfferFlow AI: Intelligent quote and invoice automation for service businesses
• AccountingML: Machine learning models for financial data analysis and prediction
• OCR Converter: Advanced optical character recognition for document digitization

Each project demonstrates our commitment to practical AI solutions that solve real-world business challenges.""",
        "Python, ML, OCR, SaaS, Automation, NLP",
        "PORTFOLIO HIGHLIGHTS"
    )
    
    # Slide 8: Contact / About
    slide = prs.slides.add_slide(slide_layout)
    
    # Set background
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = bg_color
    
    # Header
    header_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(12.33), Inches(1))
    header_frame = header_box.text_frame
    header_frame.text = "Let's Build Something Intelligent Together"
    header_para = header_frame.paragraphs[0]
    header_para.alignment = PP_ALIGN.CENTER
    header_para.font.size = Pt(28)
    header_para.font.bold = True
    header_para.font.color.rgb = accent_color
    
    # Profile section
    profile_box = slide.shapes.add_textbox(Inches(0.5), Inches(2), Inches(12.33), Inches(1.5))
    profile_frame = profile_box.text_frame
    profile_text = """Gergo Kiss, MSc
Managing Director & AI Systems Architect
Specialized in enterprise AI solutions, machine learning, and intelligent automation"""
    profile_frame.text = profile_text
    profile_para = profile_frame.paragraphs[0]
    profile_para.alignment = PP_ALIGN.CENTER
    profile_para.font.size = Pt(16)
    profile_para.font.color.rgb = text_color
    
    for para in profile_frame.paragraphs[1:]:
        para.alignment = PP_ALIGN.CENTER
        para.font.size = Pt(14)
        para.font.color.rgb = text_color
    
    # Company addresses
    address_box = slide.shapes.add_textbox(Inches(0.5), Inches(4), Inches(12.33), Inches(2.5))
    address_frame = address_box.text_frame
    address_text = """Y2K Global OÜ
Sepapaja 6, 15551 Tallinn, Estonia
Reg: 17373395, VAT: EE102928187

KISS IT Solutions e.U.
Seitenstettengasse 5/37, 1010 Vienna, Austria
Reg: FN606200x, VAT: ATU68895724

Contact: +43 1 442 20 143 | office@y2k.global | www.y2k.global"""
    address_frame.text = address_text
    address_para = address_frame.paragraphs[0]
    address_para.alignment = PP_ALIGN.CENTER
    address_para.font.size = Pt(12)
    address_para.font.color.rgb = text_color
    
    for para in address_frame.paragraphs[1:]:
        para.alignment = PP_ALIGN.CENTER
        para.font.size = Pt(12)
        para.font.color.rgb = text_color
    
    # Save the presentation
    prs.save('/git/y2k-portfolio/Y2K-AI-Portfolio.pptx')
    print("AI Portfolio presentation created successfully!")

if __name__ == "__main__":
    create_ai_portfolio()